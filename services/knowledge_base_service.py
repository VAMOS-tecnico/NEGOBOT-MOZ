"""Tenant-scoped document ingestion for the platform assistant knowledge base.

The service deliberately stores extracted text as metadata in Firestore and keeps
original bytes outside the database. Firebase Storage is preferred when a bucket
is configured; a local directory can be used as a controlled fallback for
installations that have not configured Firebase Storage yet.
"""

from __future__ import annotations

import csv
import io
import logging
import os
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from werkzeug.utils import secure_filename

logger = logging.getLogger(__name__)

MAX_FILE_BYTES = max(1, int(os.getenv("KNOWLEDGE_MAX_FILE_BYTES", str(16 * 1024 * 1024))))
MAX_EXTRACTED_CHARS = max(1000, int(os.getenv("KNOWLEDGE_MAX_EXTRACTED_CHARS", "120000")))
LOCAL_STORAGE_DIR = Path(os.getenv("KNOWLEDGE_UPLOAD_DIR", "/tmp/negobot-knowledge")).resolve()

ALLOWED_EXTENSIONS = {".pdf", ".xlsx", ".csv", ".pptx", ".docx", ".png", ".jpg", ".jpeg"}
EXTENSION_MIME_TYPES = {
    ".pdf": {"application/pdf"},
    ".xlsx": {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
    ".csv": {"text/csv", "text/plain", "application/csv", "application/vnd.ms-excel"},
    ".pptx": {"application/vnd.openxmlformats-officedocument.presentationml.presentation"},
    ".docx": {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
    ".png": {"image/png"},
    ".jpg": {"image/jpeg"},
    ".jpeg": {"image/jpeg"},
}


class KnowledgeBaseError(ValueError):
    """Expected validation or document parsing failure safe to show to a user."""


def _now() -> datetime:
    return datetime.now(timezone.utc)


def file_extension(filename: str) -> str:
    return Path(str(filename or "")).suffix.casefold()


def validate_upload(filename: str, mime_type: str, size_bytes: int) -> tuple[str, str]:
    """Validate extension, MIME and size and return a safe name plus extension."""
    original = str(filename or "").strip()
    safe_name = secure_filename(original)[:180]
    extension = file_extension(safe_name or original)
    if not safe_name or extension not in ALLOWED_EXTENSIONS:
        raise KnowledgeBaseError("Formato não suportado. Envia PDF, XLSX, CSV, PPTX, DOCX, PNG ou JPG.")
    if size_bytes > MAX_FILE_BYTES:
        raise KnowledgeBaseError(f"O ficheiro não pode exceder {MAX_FILE_BYTES // (1024 * 1024)} MB.")
    supplied_mime = str(mime_type or "").split(";", 1)[0].strip().casefold()
    allowed_mimes = EXTENSION_MIME_TYPES[extension]
    # Browsers sometimes send application/octet-stream for valid local files;
    # extension remains mandatory and MIME is rejected only when clearly wrong.
    if supplied_mime and supplied_mime != "application/octet-stream" and supplied_mime not in allowed_mimes:
        raise KnowledgeBaseError("O tipo MIME não corresponde à extensão do ficheiro.")
    return safe_name, extension


def _normalise_text(value: Any) -> str:
    text = str(value or "").replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _limit_text(text: str) -> str:
    normalised = _normalise_text(text)
    if len(normalised) <= MAX_EXTRACTED_CHARS:
        return normalised
    suffix = "\n\n[Conteúdo adicional omitido por limite de segurança.]"
    return normalised[: max(0, MAX_EXTRACTED_CHARS - len(suffix))].rstrip() + suffix


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = _normalise_text(page.extract_text() or "")
        if text:
            pages.append(f"--- PÁGINA {index} ---\n{text}")
    return "\n\n".join(pages)


def _extract_csv(content: bytes) -> str:
    decoded = ""
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            decoded = content.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if not decoded:
        raise KnowledgeBaseError("Não foi possível ler a codificação do CSV.")
    sample = decoded[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel
    rows = csv.reader(io.StringIO(decoded), dialect)
    lines = ["\t".join(_normalise_text(cell) for cell in row) for row in rows]
    return "\n".join(line for line in lines if line.strip())


def _extract_xlsx(content: bytes) -> str:
    import pandas as pd

    try:
        sheets = pd.read_excel(io.BytesIO(content), sheet_name=None, engine="openpyxl")
    except Exception as exc:
        raise KnowledgeBaseError("Não foi possível ler a folha Excel. Confirma se o ficheiro não está corrompido.") from exc
    blocks = []
    for sheet_name, frame in sheets.items():
        frame = frame.fillna("")
        rows = ["\t".join(_normalise_text(value) for value in frame.columns)]
        rows.extend("\t".join(_normalise_text(value) for value in row) for row in frame.astype(str).itertuples(index=False, name=None))
        blocks.append(f"--- FOLHA EXCEL: {sheet_name} ---\n" + "\n".join(rows))
    return "\n\n".join(blocks)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise KnowledgeBaseError("Não foi possível ler a apresentação PowerPoint.") from exc
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append("\t".join(cell.text for cell in row.cells))
        text = _normalise_text("\n".join(parts))
        if text:
            slides.append(f"--- SLIDE {index} ---\n{text}")
    return "\n\n".join(slides)


def _extract_docx(content: bytes) -> str:
    from docx import Document

    try:
        document = Document(io.BytesIO(content))
    except Exception as exc:
        raise KnowledgeBaseError("Não foi possível ler o documento Word.") from exc
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table_index, table in enumerate(document.tables, start=1):
        parts.append(f"--- TABELA {table_index} ---")
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_image(content: bytes) -> str:
    from PIL import Image

    try:
        image = Image.open(io.BytesIO(content))
        image.verify()
        image = Image.open(io.BytesIO(content))
    except Exception as exc:
        raise KnowledgeBaseError("A imagem não é válida ou está corrompida.") from exc
    try:
        import pytesseract

        language = os.getenv("KNOWLEDGE_OCR_LANG", "por+eng")
        text = pytesseract.image_to_string(image, lang=language)
    except Exception as exc:
        logger.warning("OCR indisponível para base de conhecimento: %s", type(exc).__name__)
        raise KnowledgeBaseError("Não foi possível executar OCR nesta imagem. Tenta uma imagem mais nítida.") from exc
    return text


def extract_text(filename: str, content: bytes) -> str:
    """Extract text/tables from one validated file."""
    extension = file_extension(filename)
    extractors = {
        ".pdf": _extract_pdf,
        ".csv": _extract_csv,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
        ".docx": _extract_docx,
        ".png": _extract_image,
        ".jpg": _extract_image,
        ".jpeg": _extract_image,
    }
    extractor = extractors.get(extension)
    if extractor is None:
        raise KnowledgeBaseError("Formato não suportado.")
    text = _limit_text(extractor(content))
    if not text:
        raise KnowledgeBaseError("Não foi encontrado texto legível neste ficheiro.")
    return text


def _firebase_bucket():
    try:
        from firebase_admin import storage

        bucket_name = str(os.getenv("FIREBASE_STORAGE_BUCKET") or "").strip() or None
        return storage.bucket(bucket_name) if bucket_name else storage.bucket()
    except Exception:
        return None


def store_blob(tenant_id: str, file_id: str, filename: str, content: bytes, mime_type: str, prefix: str = "knowledge") -> str | None:
    """Store bytes outside Firestore and return an opaque storage key."""
    safe_name = secure_filename(filename)[:180] or "documento"
    safe_prefix = secure_filename(prefix)[:40] or "uploads"
    key = f"{safe_prefix}/{secure_filename(tenant_id)[:80]}/{file_id}/{safe_name}"
    bucket = _firebase_bucket()
    if bucket is not None:
        try:
            blob = bucket.blob(key)
            blob.upload_from_string(content, content_type=mime_type or "application/octet-stream")
            return key
        except Exception:
            logger.exception("Falha ao guardar ficheiro no Firebase Storage")
    try:
        relative_key = f"{safe_prefix}/{secure_filename(tenant_id)[:80]}/{file_id}/{safe_name}"
        target = (LOCAL_STORAGE_DIR / relative_key).resolve()
        if LOCAL_STORAGE_DIR not in target.parents:
            raise KnowledgeBaseError("Caminho de armazenamento inválido.")
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_bytes(content)
        return f"local:{relative_key}"
    except Exception:
        logger.exception("Falha ao guardar ficheiro no armazenamento local")
        return None


def store_original(tenant_id: str, file_id: str, filename: str, content: bytes, mime_type: str) -> str | None:
    return store_blob(tenant_id, file_id, filename, content, mime_type, prefix="knowledge")


def read_blob(storage_key: str | None) -> bytes | None:
    if not storage_key:
        return None
    if storage_key.startswith("local:"):
        try:
            relative_key = storage_key.removeprefix("local:")
            path = (LOCAL_STORAGE_DIR / relative_key).resolve()
            if LOCAL_STORAGE_DIR in path.parents and path.is_file():
                return path.read_bytes()
        except Exception:
            logger.warning("Falha ao ler blob local", exc_info=True)
        return None
    bucket = _firebase_bucket()
    if bucket is None:
        return None
    try:
        return bucket.blob(storage_key).download_as_bytes()
    except Exception:
        logger.warning("Falha ao ler blob do Firebase Storage", exc_info=True)
        return None


def delete_original(storage_key: str | None) -> None:
    if not storage_key:
        return
    if storage_key.startswith("local:"):
        try:
            relative_key = storage_key.removeprefix("local:")
            path = (LOCAL_STORAGE_DIR / relative_key).resolve()
            if LOCAL_STORAGE_DIR in path.parents and path.is_file():
                path.unlink()
        except Exception:
            logger.warning("Falha ao remover cópia local da base de conhecimento", exc_info=True)
        return
    bucket = _firebase_bucket()
    if bucket is not None:
        try:
            bucket.blob(storage_key).delete()
        except Exception:
            logger.warning("Falha ao remover objecto do Firebase Storage", exc_info=True)


def serialise_file(document_id: str, data: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": document_id,
        "file_name": str(data.get("file_name") or "documento"),
        "size_bytes": int(data.get("size_bytes") or 0),
        "mime_type": str(data.get("mime_type") or "application/octet-stream"),
        "extension": str(data.get("extension") or file_extension(data.get("file_name") or "")),
        "status": str(data.get("status") or "error"),
        "error": str(data.get("error") or "") or None,
        "extracted_chars": int(data.get("extracted_chars") or 0),
        "created_at": str(data.get("created_at") or "") or None,
        "updated_at": str(data.get("updated_at") or "") or None,
        "indexed_at": str(data.get("indexed_at") or "") or None,
    }


def list_tenant_files(db, tenant_id: str, limit: int = 100) -> list[dict[str, Any]]:
    rows = []
    documents = db.collection("assistant_knowledge_files").where("tenant_id", "==", tenant_id).limit(limit).stream()
    for document in documents:
        rows.append(serialise_file(document.id, document.to_dict() or {}))
    rows.sort(key=lambda item: item.get("created_at") or "", reverse=True)
    return rows


def build_tenant_context(db, tenant_id: str, max_chars: int = 6500) -> str:
    """Return bounded, indexed-only context for the tenant's AI prompt."""
    if not tenant_id:
        return ""
    try:
        documents = db.collection("assistant_knowledge_files").where("tenant_id", "==", tenant_id).limit(100).stream()
    except Exception:
        logger.exception("Falha ao carregar contexto da base de conhecimento tenant=%s", tenant_id)
        return ""
    blocks = []
    used = 0
    for document in documents:
        data = document.to_dict() or {}
        if data.get("status") != "indexed":
            continue
        name = str(data.get("file_name") or "documento")
        text = _normalise_text(data.get("extracted_text") or "")
        if not text:
            continue
        block = f"--- FICHEIRO: {name} ---\n{text}"
        remaining = max_chars - used
        if remaining <= 0:
            break
        blocks.append(block[:remaining])
        used += min(len(block), remaining)
    return "\n\n".join(blocks)
