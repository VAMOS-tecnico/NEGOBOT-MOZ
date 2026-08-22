import io
import os
import sys
import types
import unittest
from unittest.mock import patch

sys.modules.setdefault("edge_tts", types.ModuleType("edge_tts"))
os.environ.setdefault("PLATFORM_SECRET_KEY", "platform-test-secret")
os.environ.setdefault("ADMIN_TOKEN", "admin-test-token")
os.environ.setdefault("TELEGRAM_TOKEN_ENCRYPTION_KEY", "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA=")

import extensions
extensions.init_extensions = lambda _app: None
from app import app
from services.knowledge_base_service import build_tenant_context, extract_text


class FakeSnapshot:
    def __init__(self, collection, document_id, data=None, exists=True):
        self.collection_ref = collection
        self.id = document_id
        self.data = dict(data or {})
        self.exists = exists

    def to_dict(self):
        return dict(self.data)

    def get(self):
        return self

    def set(self, values, merge=False):
        self.data = {**self.data, **values} if merge else dict(values)
        self.exists = True

    def delete(self):
        self.collection_ref.documents.pop(self.id, None)
        self.exists = False


class FakeCollection:
    def __init__(self, db, name, filters=None):
        self.db = db
        self.name = name
        self.filters = filters or []
        self.documents = db.collections.setdefault(name, {})

    def where(self, field, _operator, value):
        return FakeCollection(self.db, self.name, self.filters + [(field, value)])

    def limit(self, _amount):
        return self

    def stream(self):
        rows = list(self.documents.values())
        for field, expected in self.filters:
            rows = [row for row in rows if row.data.get(field) == expected]
        return rows

    def document(self, document_id):
        if document_id not in self.documents:
            self.documents[document_id] = FakeSnapshot(self, document_id, {}, exists=False)
        return self.documents[document_id]


class FakeDB:
    def __init__(self):
        self.collections = {}

    def collection(self, name):
        return FakeCollection(self, name)


def set_identity(client, tenant_id, tenant_role="owner"):
    with client.session_transaction() as session:
        session["platform_identity"] = {
            "id": f"user-{tenant_id}",
            "name": "Test User",
            "role": "client",
            "tenant_id": tenant_id,
            "tenant_role": tenant_role,
        }


class AssistantKnowledgeTests(unittest.TestCase):
    def setUp(self):
        app.config.update(TESTING=True, SECRET_KEY="platform-test-secret")
        self.db = FakeDB()
        extensions.db = self.db
        self.client = app.test_client()
        set_identity(self.client, "tenant-a")

    def test_extract_csv_tables(self):
        text = extract_text("catalogo.csv", "Produto;Preço\nCaderno;100\n".encode("utf-8"))
        self.assertIn("Produto", text)
        self.assertIn("Caderno", text)
        self.assertIn("100", text)

    def test_extract_docx_paragraphs_and_tables(self):
        from docx import Document

        document = Document()
        document.add_paragraph("Política de devoluções")
        table = document.add_table(rows=1, cols=2)
        table.rows[0].cells[0].text = "Prazo"
        table.rows[0].cells[1].text = "30 dias"
        content = io.BytesIO()
        document.save(content)
        text = extract_text("politica.docx", content.getvalue())
        self.assertIn("Política de devoluções", text)
        self.assertIn("30 dias", text)

    def test_extract_pptx_slides(self):
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        textbox.text = "Apresentação de produto"
        content = io.BytesIO()
        presentation.save(content)
        text = extract_text("produto.pptx", content.getvalue())
        self.assertIn("SLIDE 1", text)
        self.assertIn("Apresentação de produto", text)

    def test_upload_indexes_file_without_persisting_binary_in_firestore(self):
        with patch("routes.platform_routes.store_original", return_value="/tmp/test-knowledge/catalogo.csv"), patch("routes.platform_routes.extract_text", return_value="Produto Caderno; preço 100 MT"):
            response = self.client.post(
                "/api/platform/client/assistant/knowledge",
                data={"file": (io.BytesIO(b"Produto,Preco\nCaderno,100\n"), "catalogo.csv")},
                content_type="multipart/form-data",
            )
        self.assertEqual(response.status_code, 201)
        payload = response.get_json()
        self.assertEqual(payload["file"]["status"], "indexed")
        stored = next(iter(self.db.collections["assistant_knowledge_files"].values())).data
        self.assertEqual(stored["tenant_id"], "tenant-a")
        self.assertEqual(stored["extracted_text"], "Produto Caderno; preço 100 MT")
        self.assertNotIn("Produto,Preco", repr(stored))

    def test_invalid_extension_is_rejected(self):
        response = self.client.post(
            "/api/platform/client/assistant/knowledge",
            data={"file": (io.BytesIO(b"secret"), "secrets.txt")},
            content_type="multipart/form-data",
        )
        self.assertEqual(response.status_code, 400)
        self.assertIn("Formato não suportado", response.get_json()["error"])

    def test_list_and_context_are_tenant_scoped(self):
        collection = self.db.collection("assistant_knowledge_files")
        collection.document("file-a-123456").set({"tenant_id": "tenant-a", "file_name": "a.pdf", "status": "indexed", "extracted_text": "Oferta tenant A", "size_bytes": 10})
        collection.document("file-b-123456").set({"tenant_id": "tenant-b", "file_name": "b.pdf", "status": "indexed", "extracted_text": "Segredo tenant B", "size_bytes": 10})
        response = self.client.get("/api/platform/client/assistant/knowledge")
        self.assertEqual(response.status_code, 200)
        self.assertEqual([item["file_name"] for item in response.get_json()["files"]], ["a.pdf"])
        context = build_tenant_context(self.db, "tenant-a")
        self.assertIn("Oferta tenant A", context)
        self.assertNotIn("Segredo tenant B", context)

    def test_delete_cannot_cross_tenant(self):
        collection = self.db.collection("assistant_knowledge_files")
        collection.document("file-b-123456").set({"tenant_id": "tenant-b", "file_name": "b.pdf", "status": "indexed"})
        response = self.client.delete("/api/platform/client/assistant/knowledge/file-b-123456")
        self.assertEqual(response.status_code, 404)
        self.assertIn("não encontrado neste tenant", response.get_json()["error"])
        self.assertIn("file-b-123456", collection.documents)


if __name__ == "__main__":
    unittest.main()
